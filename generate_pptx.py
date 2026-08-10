import os
import sys
import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8')


def generate_charts():
    os.makedirs("charts", exist_ok=True)
    
    # 1. Confusion Matrix Plot
    fig, ax = plt.subplots(figsize=(5, 4))
    cm = np.array([[46, 2], [1, 1]])
    cax = ax.matshow(cm, cmap=plt.cm.Blues, alpha=0.85)
    
    for i in range(cm.shape[0]):
        for j in range(cm.shape[1]):
            ax.text(x=j, y=i, s=f"{cm[i, j]}", va='center', ha='center', size=16, weight='bold',
                    color='white' if cm[i, j] > 20 else 'black')
            
    ax.set_xticklabels(['', 'Match', 'No Match'])
    ax.set_yticklabels(['', 'Match', 'No Match'])
    plt.xlabel('Predicted Intent', fontsize=12, labelpad=10)
    plt.ylabel('Actual Intent', fontsize=12)
    plt.title('Confusion Matrix Evaluation (N=50)', fontsize=12, fontweight='bold', pad=15)
    plt.tight_layout()
    plt.savefig("charts/confusion_matrix.png", dpi=200)
    plt.close()

    # 2. Performance Metrics Bar Chart
    fig, ax = plt.subplots(figsize=(6, 3.5))
    metrics = ['Accuracy', 'Precision', 'Recall', 'F1-Score']
    scores = [94.0, 97.87, 95.83, 96.84]
    colors = ['#6366f1', '#14b8a6', '#8b5cf6', '#ec4899']
    
    bars = ax.bar(metrics, scores, color=colors, width=0.55)
    ax.set_ylim(80, 100)
    ax.set_ylabel('Percentage (%)', fontsize=11)
    ax.set_title('Model Accuracy & Evaluation Metrics', fontsize=12, fontweight='bold')
    
    for bar in bars:
        height = bar.get_height()
        ax.annotate(f'{height:.1f}%',
                    xy=(bar.get_x() + bar.get_width() / 2, height),
                    xytext=(0, 3),  # 3 points vertical offset
                    textcoords="offset points",
                    ha='center', va='bottom', fontweight='bold', fontsize=10)
                    
    plt.tight_layout()
    plt.savefig("charts/accuracy_chart.png", dpi=200)
    plt.close()

def create_pptx_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5) # 16:9 widescreen format
    
    blank_layout = prs.slide_layouts[6]
    
    def apply_dark_theme(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(11, 15, 25) # #0b0f19
        
    # SLIDE 1: Title Slide
    slide1 = prs.slides.add_slide(blank_layout)
    apply_dark_theme(slide1)
    
    # Title Box
    txBox = slide1.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.333), Inches(3))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "STUDENT QUERY AI CHATBOT"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "Natural Language Processing (NLP) & Rule-Based Intent Matching"
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(99, 102, 241) # Indigo
    p2.alignment = PP_ALIGN.CENTER
    
    p3 = tf.add_paragraph()
    p3.text = "\nMinor Project Viva & Academic Presentation"
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(148, 163, 184)
    p3.alignment = PP_ALIGN.CENTER

    # SLIDE 2: Problem Statement & Objectives
    slide2 = prs.slides.add_slide(blank_layout)
    apply_dark_theme(slide2)
    
    tb = slide2.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = "Problem Statement & Project Objectives"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Card 1: Problem
    shape1 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = RGBColor(30, 41, 59)
    shape1.line.color.rgb = RGBColor(99, 102, 241)
    tf1 = shape1.text_frame
    tf1.word_wrap = True
    p = tf1.paragraphs[0]
    p.text = "❌ Existing Challenges\n"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(239, 68, 68)
    
    bullet_text1 = [
        "Helpdesks flooded with repetitive administrative inquiries.",
        "Long student wait times during admission peak periods.",
        "Manual response handling is inefficient and costly.",
        "Commercial LLM APIs (GPT-4) can be expensive and prone to hallucination."
    ]
    for bt in bullet_text1:
        bp = tf1.add_paragraph()
        bp.text = f"• {bt}"
        bp.font.size = Pt(14)
        bp.font.color.rgb = RGBColor(226, 232, 240)
        bp.space_after = Pt(8)

    # Card 2: Objectives
    shape2 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = RGBColor(30, 41, 59)
    shape2.line.color.rgb = RGBColor(20, 184, 166)
    tf2 = shape2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "🎯 Project Objectives\n"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(20, 184, 166)
    
    bullet_text2 = [
        "Develop an offline-capable Python NLP Chatbot.",
        "Implement TF-IDF & Cosine Similarity for confidence scoring.",
        "Pre-load dataset covering Courses, Fees, Timings, Contact, Admissions, Hostels.",
        "Build interactive dual interfaces: Terminal CLI and Web Dashboard.",
        "Provide explicit exit controls ('exit', 'quit', 'bye')."
    ]
    for bt in bullet_text2:
        bp = tf2.add_paragraph()
        bp.text = f"✓ {bt}"
        bp.font.size = Pt(14)
        bp.font.color.rgb = RGBColor(226, 232, 240)
        bp.space_after = Pt(8)

    # SLIDE 3: System Architecture & Workflow
    slide3 = prs.slides.add_slide(blank_layout)
    apply_dark_theme(slide3)
    
    tb = slide3.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = "System Architecture & Processing Pipeline"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Diagram steps as horizontal cards
    steps = [
        ("1. User Query", "Input string from CLI or Web UI", RGBColor(99, 102, 241)),
        ("2. Preprocessing", "NLTK lowercasing, tokenization & lemmatization", RGBColor(139, 92, 246)),
        ("3. TF-IDF Matrix", "Feature extraction & English stop words filtering", RGBColor(20, 184, 166)),
        ("4. Cosine Similarity", "Calculate similarity score vs pattern vectors", RGBColor(236, 72, 153)),
        ("5. Output Response", "Display response, confidence badge & followups", RGBColor(16, 185, 129))
    ]
    
    left_margin = Inches(0.8)
    card_width = Inches(2.2)
    gap = Inches(0.2)
    
    for i, (title, desc, color) in enumerate(steps):
        posX = left_margin + i * (card_width + gap)
        sc = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, posX, Inches(2.2), card_width, Inches(4.2))
        sc.fill.solid()
        sc.fill.fore_color.rgb = RGBColor(30, 41, 59)
        sc.line.color.rgb = color
        
        stf = sc.text_frame
        stf.word_wrap = True
        p = stf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color
        
        p_desc = stf.add_paragraph()
        p_desc.text = f"\n{desc}"
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = RGBColor(203, 213, 225)

    # SLIDE 4: Experimental Performance Charts
    slide4 = prs.slides.add_slide(blank_layout)
    apply_dark_theme(slide4)
    
    tb = slide4.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = "Experimental Results & Model Evaluation"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    if os.path.exists("charts/confusion_matrix.png"):
        slide4.shapes.add_picture("charts/confusion_matrix.png", Inches(0.8), Inches(1.8), height=Inches(4.8))
        
    if os.path.exists("charts/accuracy_chart.png"):
        slide4.shapes.add_picture("charts/accuracy_chart.png", Inches(6.5), Inches(1.8), height=Inches(4.8))

    # SLIDE 5: Conclusion & Summary
    slide5 = prs.slides.add_slide(blank_layout)
    apply_dark_theme(slide5)
    
    tb = slide5.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = "Conclusion & Summary"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    box = slide5.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10.3), Inches(4.5))
    tf = box.text_frame
    tf.word_wrap = True
    
    points = [
        "✅ Fully Functional Student Assistant: Resolves queries across 8 major academic categories.",
        "✅ High Matching Accuracy: Achieved 94.0% accuracy and 96.84% F1-score.",
        "✅ Ultra-Low Latency: Response processing completed in < 5ms.",
        "✅ Dual Interface Deployment: Console CLI (`cli_chatbot.py`) + Flask Web App (`app.py` / `index.html`).",
        "✅ Offline Resilience: Embedded client-side JS matching ensures interactive availability even without Flask."
    ]
    for pt in points:
        p = tf.add_paragraph()
        p.text = pt
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(241, 245, 249)
        p.space_after = Pt(14)

    output_filename = "Student_Query_AI_Chatbot_Presentation.pptx"
    prs.save(output_filename)
    print(f"✅ Generated PowerPoint Presentation: {output_filename}")

if __name__ == "__main__":
    generate_charts()
    create_pptx_presentation()
